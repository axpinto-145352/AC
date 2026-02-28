#!/usr/bin/env python3
"""
Generate DIU Solution Brief PowerPoint: AI-Assisted Triage & Treatment (PROJ00628)
Veteran Vectors (Prime) | Authentic Consortium
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
GOLD   = RGBColor(0xC5, 0x96, 0x0C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
MGRAY  = RGBColor(0xE0, 0xE0, 0xE0)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
LTNAVY = RGBColor(0x2C, 0x3E, 0x6B)
PALEBLUE = RGBColor(0xE8, 0xEE, 0xF6)

FONT_NAME = "Calibri"

# Slide dimensions (standard widescreen 13.333 x 7.5 in)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
BLANK_LAYOUT = prs.slide_layouts[6]


# ===========================================================================
# Helper functions
# ===========================================================================

def add_solid_bg(slide, color):
    """Set slide background to a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=14, bold=False, color=DGRAY,
             alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tf


def add_paragraph(tf, text, font_size=14, bold=False, color=DGRAY,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4),
                  space_after=Pt(2), font_name=FONT_NAME, level=0):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=DGRAY, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_header_bar(slide, title_text):
    """Add a navy header bar at top with gold accent line and slide title."""
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    accent = add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.06), GOLD)
    add_textbox(slide, Inches(0.6), Inches(0.18), Inches(10), Inches(0.8),
                title_text, font_size=28, bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT)
    # Small branding on right
    add_textbox(slide, Inches(9.5), Inches(0.25), Inches(3.5), Inches(0.6),
                "Veteran Vectors  |  Authentic Consortium",
                font_size=10, bold=False, color=RGBColor(0xA0, 0xB0, 0xC8),
                alignment=PP_ALIGN.RIGHT)


def add_footer(slide, slide_num, total=15):
    """Small footer with PROJ00628 and slide number."""
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(5), Inches(0.35),
                "PROJ00628  |  AI-Assisted Triage & Treatment  |  UNCLASSIFIED",
                font_size=8, color=RGBColor(0x99, 0x99, 0x99))
    add_textbox(slide, Inches(10), Inches(7.05), Inches(2.8), Inches(0.35),
                f"{slide_num} / {total}",
                font_size=8, color=RGBColor(0x99, 0x99, 0x99),
                alignment=PP_ALIGN.RIGHT)


def bullet_textbox(slide, left, top, width, height, bullets,
                   font_size=14, color=DGRAY, bold_first_word=False,
                   bullet_char="\u2022", spacing=Pt(6)):
    """Create a textbox with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        p.space_after = Pt(2)
        p.alignment = PP_ALIGN.LEFT
        if bold_first_word and ":" in item:
            colon_idx = item.index(":")
            prefix = item[:colon_idx + 1]
            suffix = item[colon_idx + 1:]
            run1 = p.add_run()
            run1.text = f"{bullet_char}  {prefix}"
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = FONT_NAME
            run2 = p.add_run()
            run2.text = suffix
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = FONT_NAME
        else:
            run = p.add_run()
            run.text = f"{bullet_char}  {item}"
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = color
            run.font.name = FONT_NAME
    return txBox


def add_icon_box(slide, left, top, width, height, icon_text, label,
                 fill=PALEBLUE, icon_color=NAVY, label_color=DGRAY):
    """Rounded rectangle with icon/emoji text on top, label below."""
    box = add_rounded_rect(slide, left, top, width, height,
                           fill_color=fill, line_color=MGRAY, line_width=Pt(1))
    # Icon text (large)
    add_textbox(slide, left, top + Inches(0.1), width, Inches(0.5),
                icon_text, font_size=22, bold=True, color=icon_color,
                alignment=PP_ALIGN.CENTER)
    # Label (below)
    add_textbox(slide, left, top + Inches(0.55), width, Inches(0.8),
                label, font_size=11, bold=False, color=label_color,
                alignment=PP_ALIGN.CENTER)
    return box


def add_stat_box(slide, left, top, width, height, number, label,
                 num_color=GOLD, label_color=DGRAY, fill=None):
    box = add_rounded_rect(slide, left, top, width, height,
                           fill_color=fill or PALEBLUE,
                           line_color=MGRAY, line_width=Pt(1))
    add_textbox(slide, left, top + Inches(0.15), width, Inches(0.6),
                number, font_size=28, bold=True, color=num_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.7), width, Inches(0.5),
                label, font_size=11, bold=False, color=label_color,
                alignment=PP_ALIGN.CENTER)
    return box


# ===========================================================================
# SLIDE 1  --  Title Slide
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, NAVY)

# Large gold accent line
add_rect(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
            "AI-Assisted Triage & Treatment Solution",
            font_size=40, bold=True, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.6),
            "VitalEdge AI  \u2014  Hemodynamic Monitoring & Decision Support for the Battlefield",
            font_size=20, bold=False, color=RGBColor(0xC0, 0xCC, 0xDD))

# Separator
add_rect(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.03), GOLD)

add_textbox(slide, Inches(0.8), Inches(3.9), Inches(5), Inches(0.5),
            "Veteran Vectors  (Prime)", font_size=22, bold=True, color=GOLD)
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
            "Authentic Consortium", font_size=18, bold=False, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(5.3), Inches(5), Inches(0.5),
            "PROJ00628  |  DIU AI-Assisted Triage & Treatment Challenge",
            font_size=14, bold=False, color=RGBColor(0x88, 0x99, 0xBB))
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(5), Inches(0.5),
            "February 2026  |  Solution Brief",
            font_size=14, bold=False, color=RGBColor(0x88, 0x99, 0xBB))

# Classification banner bottom
add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), RGBColor(0x12, 0x1E, 0x36))
add_textbox(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35),
            "UNCLASSIFIED  |  SDVOSB  |  CAGE: XXXXX",
            font_size=9, color=RGBColor(0x70, 0x80, 0x99),
            alignment=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 2  --  Problem Statement
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Problem Statement")
add_footer(slide, 2)

# Three problem cards
card_w = Inches(3.6)
card_h = Inches(3.5)
card_y = Inches(1.7)
gap = Inches(0.5)
start_x = Inches(0.8)

problems = [
    ("No Real-Time Decision Support",
     "Forward combat medics operating in DDIL (Denied, Disrupted, Intermittent, Limited) "
     "environments have zero hemodynamic decision-support tools. They rely on training and "
     "instinct under extreme stress.",
     "\u26A0"),
    ("Subjective Triage Assessment",
     "Current triage protocols depend on subjective visual and manual assessment. "
     "Cognitive overload in mass casualty events leads to mis-triage rates of 20\u201350%, "
     "directly impacting survivability.",
     "\u2139"),
    ("Communication Delays Cost Lives",
     "Casualty status updates from point-of-injury to Role 2/3 facilities are delayed "
     "or lost. Receiving medical teams lack predictive data, reducing the Golden Hour advantage.",
     "\u23F1"),
]

for i, (title, desc, icon) in enumerate(problems):
    x = start_x + i * (card_w + gap)
    card = add_rounded_rect(slide, x, card_y, card_w, card_h,
                            fill_color=PALEBLUE, line_color=MGRAY, line_width=Pt(1.5))
    # Gold top accent
    add_rect(slide, x, card_y, card_w, Inches(0.06), GOLD)
    # Icon
    add_textbox(slide, x, card_y + Inches(0.2), card_w, Inches(0.5),
                icon, font_size=28, bold=True, color=NAVY,
                alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.25), card_y + Inches(0.7), card_w - Inches(0.5), Inches(0.5),
                title, font_size=16, bold=True, color=NAVY,
                alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.25), card_y + Inches(1.3), card_w - Inches(0.5), Inches(2),
                desc, font_size=12, color=DGRAY, alignment=PP_ALIGN.LEFT)

# Bottom stat bar
add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.0),
         fill_color=RGBColor(0xF8, 0xF4, 0xE8), line_color=GOLD, line_width=Pt(1))
add_textbox(slide, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.8),
            "\"Preventable death on the battlefield remains the #1 DoD medical priority. "
            "AI-enabled triage can close the gap between injury and intervention.\"",
            font_size=13, bold=False, color=NAVY, alignment=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 3  --  Our Solution - Overview
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Our Solution \u2014 VitalEdge AI Overview")
add_footer(slide, 3)

# Tagline
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "An AI-powered hemodynamic monitoring and triage decision-support platform "
            "purpose-built for DDIL environments.",
            font_size=15, bold=False, color=NAVY, alignment=PP_ALIGN.LEFT)

# Three-tier visual
tier_w = Inches(3.2)
tier_h = Inches(3.0)
tier_y = Inches(2.2)
arrow_w = Inches(0.7)

tiers = [
    ("DEVICE TIER", "Wearable Sensors",
     ["COTS ruggedized hemodynamic sensors",
      "Heart rate, SpO2, blood pressure, temp",
      "BLE / Wi-Fi / USB connectivity",
      "72-hour battery | MIL-STD-810H"]),
    ("EDGE TIER", "AI Inference & Triage",
     ["On-device ML inference (<50ms)",
      "TCCC MARCH triage scoring",
      "Deterioration early warning",
      "Store-and-forward for DDIL ops"]),
    ("CLOUD TIER", "Analytics & EHR Sync",
     ["AWS GovCloud (IL-5)",
      "Population health analytics",
      "MHS GENESIS / TMIP integration",
      "Fleet management & OTA updates"]),
]

colors_tier = [RGBColor(0x1B, 0x2A, 0x4A), RGBColor(0x2C, 0x3E, 0x6B), RGBColor(0x3D, 0x55, 0x8C)]

for i, (tier_title, subtitle, bullets) in enumerate(tiers):
    x = Inches(0.8) + i * (tier_w + arrow_w)
    # Tier box
    box = add_rounded_rect(slide, x, tier_y, tier_w, tier_h,
                           fill_color=PALEBLUE, line_color=colors_tier[i],
                           line_width=Pt(2))
    # Tier header
    hdr = add_rect(slide, x, tier_y, tier_w, Inches(0.65), colors_tier[i])
    add_textbox(slide, x, tier_y + Inches(0.05), tier_w, Inches(0.35),
                tier_title, font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, tier_y + Inches(0.32), tier_w, Inches(0.3),
                subtitle, font_size=11, color=RGBColor(0xC0, 0xCC, 0xDD),
                alignment=PP_ALIGN.CENTER)
    # Bullets
    bullet_textbox(slide, x + Inches(0.15), tier_y + Inches(0.75),
                   tier_w - Inches(0.3), tier_h - Inches(0.9),
                   bullets, font_size=11, color=DGRAY)
    # Arrow between tiers
    if i < 2:
        ax = x + tier_w + Inches(0.08)
        add_textbox(slide, ax, tier_y + Inches(1.2), arrow_w, Inches(0.5),
                    "\u27A1", font_size=30, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)

# Bottom note
add_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.7),
         fill_color=RGBColor(0xF8, 0xF4, 0xE8), line_color=GOLD, line_width=Pt(1))
add_textbox(slide, Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.6),
            "KEY:  Operates fully disconnected  \u2022  AI recommends, medic decides (human-in-the-loop)  "
            "\u2022  Automatic sync when connectivity restored",
            font_size=12, bold=False, color=NAVY, alignment=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 4  --  System Architecture
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "System Architecture")
add_footer(slide, 4)

# Architecture table-based layout
# Device Tier column
col_w = Inches(3.5)
col_h = Inches(4.5)
row_y = Inches(1.6)

arch_cols = [
    ("DEVICE TIER", NAVY, [
        ("Hemodynamic Sensors", "HR, SpO2, BP, Temp, ECG"),
        ("Wearable Config", "Limb-mount, chest-strap"),
        ("Standoff Config", "Handheld probe"),
        ("Comms Module", "BLE 5.0, Wi-Fi 6, USB-C"),
    ]),
    ("EDGE TIER (BATDOK)", LTNAVY, [
        ("ONNX Runtime Engine", "ML model inference <50ms"),
        ("Triage Algorithm", "MARCH protocol scoring"),
        ("Early Warning", "15-30 min deterioration pred."),
        ("Local Data Store", "SQLite + store-and-forward"),
        ("ATAK Plugin", "Blue Force Tracker overlay"),
    ]),
    ("CLOUD TIER (GovCloud)", RGBColor(0x3D, 0x55, 0x8C), [
        ("API Gateway", "REST + MQTT broker"),
        ("Analytics Engine", "Population health / trends"),
        ("EHR Integration", "HL7 FHIR \u2192 MHS GENESIS"),
        ("Fleet Management", "OTA updates, device health"),
        ("Auth & Security", "CAC/PIV, AES-256, TLS 1.3"),
    ]),
]

for ci, (col_title, col_color, components) in enumerate(arch_cols):
    x = Inches(0.6) + ci * (col_w + Inches(0.55))
    # Column header
    add_rect(slide, x, row_y, col_w, Inches(0.55), col_color)
    add_textbox(slide, x, row_y + Inches(0.08), col_w, Inches(0.4),
                col_title, font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    # Component rows
    for ri, (comp, desc) in enumerate(components):
        ry = row_y + Inches(0.6) + ri * Inches(0.75)
        fill = PALEBLUE if ri % 2 == 0 else WHITE
        add_rounded_rect(slide, x + Inches(0.1), ry, col_w - Inches(0.2), Inches(0.65),
                         fill_color=fill, line_color=MGRAY, line_width=Pt(0.75))
        add_textbox(slide, x + Inches(0.25), ry + Inches(0.04),
                    col_w - Inches(0.5), Inches(0.3),
                    comp, font_size=11, bold=True, color=NAVY)
        add_textbox(slide, x + Inches(0.25), ry + Inches(0.32),
                    col_w - Inches(0.5), Inches(0.3),
                    desc, font_size=10, color=DGRAY)
    # Arrows between columns
    if ci < 2:
        ax = x + col_w + Inches(0.05)
        add_textbox(slide, ax, row_y + Inches(2.0), Inches(0.5), Inches(0.5),
                    "\u27A1", font_size=28, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)

# Data flow label
add_rect(slide, Inches(0.6), Inches(6.3), Inches(11.7), Inches(0.5),
         fill_color=RGBColor(0xF8, 0xF4, 0xE8), line_color=GOLD, line_width=Pt(1))
add_textbox(slide, Inches(0.8), Inches(6.33), Inches(11.3), Inches(0.4),
            "DATA FLOW:  Sensors \u2192 BLE \u2192 Edge (BATDOK) \u2192 Mesh/SATCOM \u2192 Cloud (GovCloud)  "
            "|  Store-and-forward ensures DDIL resilience",
            font_size=11, bold=False, color=NAVY, alignment=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 5  --  AI/ML Capability
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "AI/ML Capability")
add_footer(slide, 5)

# Left column - capabilities
cap_x = Inches(0.8)
cap_y = Inches(1.5)
cap_w = Inches(7.5)

caps = [
    "Triage Classification: Real-time TCCC MARCH protocol scoring (Massive hemorrhage, "
    "Airway, Respiration, Circulation, Hypothermia)",
    "Deterioration Early Warning: 15\u201330 minute predictive model alerts medics before "
    "clinical signs manifest",
    "Training Data: MIMIC-III/IV critical care database + DoD Trauma Registry + "
    "synthetic augmentation for combat-specific scenarios",
    "Edge Inference: <50ms latency on ARM-based edge compute \u2014 no cloud dependency required",
    "Human-in-the-Loop: AI generates recommendations and confidence scores; medic retains "
    "full decision authority",
    "Model Format: ONNX for cross-platform portability; quantized INT8 for edge efficiency",
]

bullet_textbox(slide, cap_x, cap_y, cap_w, Inches(4.0), caps,
               font_size=13, color=DGRAY, bold_first_word=True, spacing=Pt(10))

# Right column - stat boxes
stat_x = Inches(9.0)
stats = [
    ("<50ms", "Edge Inference\nLatency"),
    ("15-30 min", "Predictive\nWarning Window"),
    ("MARCH", "TCCC Protocol\nCompliant"),
    ("ONNX", "Cross-Platform\nML Runtime"),
]

for i, (num, label) in enumerate(stats):
    sy = Inches(1.5) + i * Inches(1.3)
    add_stat_box(slide, stat_x, sy, Inches(3.5), Inches(1.1),
                 num, label, num_color=NAVY, fill=PALEBLUE)


# ===========================================================================
# SLIDE 6  --  Key Integrations
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Key Integrations")
add_footer(slide, 6)

integrations = [
    ("BATDOK", "Medic Tablet Display",
     ["Primary UI for triage scores & vitals", "Ruggedized Android tablet",
      "Touch-optimized casualty workflow", "Works fully offline"]),
    ("ATAK / TAK", "Commander Awareness",
     ["Blue Force Tracker casualty overlay", "Real-time 9-line MEDEVAC data",
      "Mass casualty event visualization", "Plugin architecture via TAK SDK"]),
    ("EHR / FHIR", "Clinical Continuity",
     ["HL7 FHIR R4 standard", "MHS GENESIS integration",
      "TMIP-J compatibility", "Automated encounter creation"]),
    ("Mesh Network", "Multi-Casualty Comms",
     ["MQTT-based pub/sub messaging", "Automatic peer discovery",
      "Data deduplication & sync", "Resilient to node loss"]),
]

card_w = Inches(2.7)
card_h = Inches(4.2)
start_x = Inches(0.55)
card_y = Inches(1.5)

for i, (title, subtitle, bullets) in enumerate(integrations):
    x = start_x + i * (card_w + Inches(0.25))
    # Card background
    add_rounded_rect(slide, x, card_y, card_w, card_h,
                     fill_color=WHITE, line_color=NAVY, line_width=Pt(1.5))
    # Header
    add_rect(slide, x, card_y, card_w, Inches(0.9), NAVY)
    add_textbox(slide, x, card_y + Inches(0.08), card_w, Inches(0.4),
                title, font_size=18, bold=True, color=GOLD,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, card_y + Inches(0.45), card_w, Inches(0.4),
                subtitle, font_size=11, color=RGBColor(0xA0, 0xB0, 0xC8),
                alignment=PP_ALIGN.CENTER)
    # Bullets
    bullet_textbox(slide, x + Inches(0.15), card_y + Inches(1.0),
                   card_w - Inches(0.3), card_h - Inches(1.2),
                   bullets, font_size=11, color=DGRAY, spacing=Pt(8))


# ===========================================================================
# SLIDE 7  --  Hardware Platform
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Hardware Platform")
add_footer(slide, 7)

# Left - description
hw_bullets = [
    "COTS-Based Design: Ruggedized commercial off-the-shelf hemodynamic sensor suite "
    "minimizes development risk and accelerates fielding",
    "Dual Configuration: Wearable (continuous monitoring on limb/chest) and standoff "
    "(handheld probe for rapid assessment)",
    "72-Hour Battery Life: Extended operation without resupply; hot-swappable battery packs",
    "MIL-STD-810H Compliant: Shock, vibration, temperature, humidity, altitude, "
    "sand/dust environmental testing",
    "Multi-Protocol Connectivity: BLE 5.0, Wi-Fi 6, USB-C for maximum interoperability",
    "IP67 Rating: Waterproof and dustproof for austere environments",
]

bullet_textbox(slide, Inches(0.8), Inches(1.5), Inches(7.0), Inches(4.5),
               hw_bullets, font_size=13, color=DGRAY,
               bold_first_word=True, spacing=Pt(10))

# Right - spec boxes
specs = [
    ("72 hrs", "Battery Life"),
    ("MIL-STD\n810H", "Environmental\nCompliance"),
    ("IP67", "Ingress\nProtection"),
    ("BLE 5.0", "Primary\nConnectivity"),
    ("<100g", "Sensor\nWeight"),
]

for i, (val, label) in enumerate(specs):
    row = i // 2
    col = i % 2
    sx = Inches(8.5) + col * Inches(2.3)
    sy = Inches(1.5) + row * Inches(1.5)
    add_stat_box(slide, sx, sy, Inches(2.1), Inches(1.3),
                 val, label, num_color=NAVY, fill=PALEBLUE)


# ===========================================================================
# SLIDE 8  --  FDA & Regulatory Strategy
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "FDA & Regulatory Strategy")
add_footer(slide, 8)

# Four regulatory pillars
pillars = [
    ("510(k) Pathway", [
        "Predicate device identification",
        "Substantial equivalence demonstration",
        "Class II medical device classification",
        "Target: Q3 2026 submission",
    ]),
    ("SaMD \u2014 IEC 62304", [
        "Software as Medical Device classification",
        "Full software development lifecycle",
        "Unit, integration, system-level testing",
        "Traceability matrix maintained",
    ]),
    ("ISO 14971 Risk Mgmt", [
        "Hazard analysis & risk assessment",
        "Risk control measures documented",
        "Residual risk evaluation",
        "Post-market surveillance plan",
    ]),
    ("Clinical Validation", [
        "DoD Trauma Registry benchmarking",
        "MIMIC-III/IV retrospective study",
        "Simulated combat casualty trials",
        "Sensitivity/specificity reporting",
    ]),
]

pill_w = Inches(2.7)
pill_h = Inches(4.0)
pill_y = Inches(1.6)
start_x = Inches(0.55)

for i, (title, bullets) in enumerate(pillars):
    x = start_x + i * (pill_w + Inches(0.25))
    add_rounded_rect(slide, x, pill_y, pill_w, pill_h,
                     fill_color=WHITE, line_color=NAVY, line_width=Pt(1.5))
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(1.05), pill_y - Inches(0.2),
                                  Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = GOLD
    circ.line.fill.background()
    circ.shadow.inherit = False
    set_text(circ, str(i + 1), font_size=20, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    circ.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_textbox(slide, x + Inches(0.1), pill_y + Inches(0.35), pill_w - Inches(0.2), Inches(0.5),
                title, font_size=14, bold=True, color=NAVY,
                alignment=PP_ALIGN.CENTER)
    # Bullets
    bullet_textbox(slide, x + Inches(0.15), pill_y + Inches(0.9),
                   pill_w - Inches(0.3), pill_h - Inches(1.1),
                   bullets, font_size=11, color=DGRAY, spacing=Pt(8))


# ===========================================================================
# SLIDE 9  --  Cybersecurity & Compliance
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Cybersecurity & Compliance")
add_footer(slide, 9)

# Two columns
left_items = [
    ("Infrastructure & Authorization", [
        "IL-5 ATO pathway via AWS GovCloud",
        "FedRAMP High baseline alignment",
        "Continuous ATO (cATO) methodology",
        "STIG-hardened deployment configurations",
        "Infrastructure-as-code (Terraform)",
    ]),
    ("Supply Chain Security", [
        "Software Bill of Materials (SBOM)",
        "Signed builds & verified provenance",
        "Automated vulnerability scanning",
        "Container image attestation",
        "NIST SP 800-171 compliance",
    ]),
]

right_items = [
    ("Data Protection", [
        "AES-256 encryption at rest",
        "TLS 1.3 encryption in transit",
        "Field-level encryption for PII/PHI",
        "HIPAA-compliant data handling",
        "Data sovereignty controls",
    ]),
    ("Identity & Access", [
        "CAC / PIV authentication required",
        "Role-based access control (RBAC)",
        "Multi-factor authentication",
        "Audit logging & SIEM integration",
        "Zero-trust architecture principles",
    ]),
]

for col, items in enumerate([left_items, right_items]):
    base_x = Inches(0.6) + col * Inches(6.2)
    for row, (section_title, bullets) in enumerate(items):
        by = Inches(1.5) + row * Inches(2.6)
        add_rounded_rect(slide, base_x, by, Inches(5.8), Inches(2.3),
                         fill_color=WHITE, line_color=NAVY, line_width=Pt(1.5))
        add_rect(slide, base_x, by, Inches(5.8), Inches(0.5), NAVY)
        add_textbox(slide, base_x + Inches(0.2), by + Inches(0.08),
                    Inches(5.4), Inches(0.4),
                    section_title, font_size=13, bold=True, color=WHITE)
        bullet_textbox(slide, base_x + Inches(0.2), by + Inches(0.55),
                       Inches(5.4), Inches(1.6),
                       bullets, font_size=12, color=DGRAY, spacing=Pt(5))


# ===========================================================================
# SLIDE 10  --  Scalability & Manufacturing
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Scalability & Manufacturing")
add_footer(slide, 10)

# Timeline-style boxes
scale_items = [
    ("MAY 2026", "Sword Demo", "30 units delivered for operational evaluation",
     NAVY),
    ("DEC 2026", "Limited Production", "500 units for expanded testing & feedback",
     LTNAVY),
    ("MAY 2027", "Full-Rate Production", "15,000 units/year manufacturing capacity",
     RGBColor(0x3D, 0x55, 0x8C)),
]

box_w = Inches(3.5)
box_h = Inches(1.8)

for i, (date, phase, desc, color) in enumerate(scale_items):
    x = Inches(0.6) + i * (box_w + Inches(0.35))
    y = Inches(1.6)
    add_rounded_rect(slide, x, y, box_w, box_h,
                     fill_color=WHITE, line_color=color, line_width=Pt(2))
    add_rect(slide, x, y, box_w, Inches(0.55), color)
    add_textbox(slide, x, y + Inches(0.08), box_w, Inches(0.4),
                f"{date}  \u2014  {phase}", font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.65), box_w - Inches(0.4), Inches(1.0),
                desc, font_size=13, color=DGRAY, alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # Arrow
    if i < 2:
        ax = x + box_w + Inches(0.02)
        add_textbox(slide, ax, y + Inches(0.6), Inches(0.35), Inches(0.5),
                    "\u27A1", font_size=24, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)

# Key metrics row
metrics = [
    ("$85", "Per-Unit Software\nCost at Scale"),
    ("OTA", "Over-the-Air\nUpdate Capability"),
    ("24/7", "Automated Fleet\nManagement"),
    ("CI/CD", "Continuous\nImprovement Pipeline"),
]

for i, (val, label) in enumerate(metrics):
    mx = Inches(0.6) + i * Inches(3.1)
    my = Inches(3.9)
    add_stat_box(slide, mx, my, Inches(2.8), Inches(1.2),
                 val, label, num_color=GOLD, fill=PALEBLUE)

# Manufacturing details
mfg_bullets = [
    "Domestic COTS supply chain with dual-source components for resilience",
    "Automated test & calibration stations for quality assurance at scale",
    "MIL-STD-810H environmental screening integrated into production line",
    "Packaging optimized for rapid deployment kits (squad & platoon level)",
]

bullet_textbox(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5),
               mfg_bullets, font_size=12, color=DGRAY, spacing=Pt(5))


# ===========================================================================
# SLIDE 11  --  Team & Consortium
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Team & Consortium")
add_footer(slide, 11)

# Prime contractor - large card
prime_w = Inches(11.7)
prime_h = Inches(1.6)
add_rounded_rect(slide, Inches(0.8), Inches(1.5), prime_w, prime_h,
                 fill_color=NAVY, line_color=GOLD, line_width=Pt(2))
add_textbox(slide, Inches(1.0), Inches(1.6), Inches(4), Inches(0.5),
            "VETERAN VECTORS  (Prime)", font_size=20, bold=True, color=GOLD)
add_textbox(slide, Inches(1.0), Inches(2.1), Inches(4), Inches(0.4),
            "AI/ML  |  Software Development  |  Edge Compute  |  Systems Integration",
            font_size=13, color=WHITE)
add_textbox(slide, Inches(1.0), Inches(2.5), Inches(4), Inches(0.4),
            "Service-Disabled Veteran-Owned Small Business (SDVOSB)",
            font_size=11, color=RGBColor(0xA0, 0xB0, 0xC8))
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.3), Inches(1.2),
            "Founded by combat veterans who understand the operational environment. "
            "Deep expertise in AI/ML, edge computing, and DoD software development. "
            "Proven delivery on SBIR/STTR and OTA contracts.",
            font_size=12, color=RGBColor(0xC0, 0xCC, 0xDD))

# Consortium partners
partners = [
    ("Medical Device OEM", "Hardware engineering\nFDA regulatory\nManufacturing QMS",
     RGBColor(0x8B, 0x4F, 0x0C)),
    ("Cloud / Cyber Partner", "IL-5 ATO expertise\nFedRAMP alignment\nSecurity operations",
     RGBColor(0x2C, 0x6B, 0x4A)),
    ("Manufacturing Partner", "Scale production\nSupply chain mgmt\nQuality assurance",
     RGBColor(0x4A, 0x2C, 0x6B)),
    ("Clinical SME Team", "TCCC expertise\nCombat medicine\nClinical validation",
     RGBColor(0x6B, 0x2C, 0x2C)),
]

card_w = Inches(2.7)
card_h = Inches(2.8)
card_y = Inches(3.5)

for i, (name, roles, accent) in enumerate(partners):
    x = Inches(0.6) + i * (card_w + Inches(0.3))
    add_rounded_rect(slide, x, card_y, card_w, card_h,
                     fill_color=WHITE, line_color=accent, line_width=Pt(2))
    add_rect(slide, x, card_y, card_w, Inches(0.06), accent)
    add_textbox(slide, x + Inches(0.15), card_y + Inches(0.2),
                card_w - Inches(0.3), Inches(0.5),
                name, font_size=14, bold=True, color=NAVY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), card_y + Inches(0.8),
                card_w - Inches(0.3), Inches(1.8),
                roles, font_size=12, color=DGRAY,
                alignment=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 12  --  Competitive Advantages
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Competitive Advantages")
add_footer(slide, 12)

advantages = [
    ("Veteran-Owned", "Founded by combat veterans with firsthand understanding of "
     "battlefield medicine challenges and operational environments"),
    ("AI-First Approach", "Purpose-built AI/ML engine vs. sensor-only incumbents; "
     "predictive analytics, not just data display"),
    ("DDIL-Native Architecture", "Designed ground-up for disconnected operations \u2014 "
     "not a cloud product with an offline mode bolted on"),
    ("Open Standards", "HL7 FHIR, MQTT, ONNX ensure interoperability and avoid vendor "
     "lock-in across the DoD ecosystem"),
    ("Consortium Agility", "Small business speed and innovation with large-enterprise "
     "manufacturing and regulatory capability"),
]

for i, (title, desc) in enumerate(advantages):
    y = Inches(1.5) + i * Inches(1.1)
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(0.8), y + Inches(0.1),
                                   Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    badge.shadow.inherit = False
    set_text(badge, str(i + 1), font_size=16, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Content
    add_textbox(slide, Inches(1.5), y, Inches(2.5), Inches(0.4),
                title, font_size=15, bold=True, color=NAVY)
    add_textbox(slide, Inches(1.5), y + Inches(0.35), Inches(10.5), Inches(0.6),
                desc, font_size=12, color=DGRAY)
    # Subtle line
    if i < 4:
        add_rect(slide, Inches(1.5), y + Inches(0.9), Inches(10.5), Inches(0.01), MGRAY)


# ===========================================================================
# SLIDE 13  --  Timeline & Milestones
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Timeline & Milestones")
add_footer(slide, 13)

# Timeline bar
bar_y = Inches(2.8)
bar_x = Inches(0.8)
bar_w = Inches(11.7)
add_rect(slide, bar_x, bar_y, bar_w, Inches(0.08), NAVY)

milestones = [
    ("MAR\n2026", "Solution Brief\nSubmitted", 0.0),
    ("APR\n2026", "Pitch Days\nPresentation", 0.2),
    ("MAY\n2026", "Sword 2026 Demo\n30 Units", 0.35),
    ("JUN-DEC\n2026", "FDA 510(k)\nIL-5 ATO", 0.6),
    ("MAY\n2027", "Full Production\n15,000 units/yr", 1.0),
]

for date, desc, pct in milestones:
    mx = bar_x + Emu(int(bar_w.emu * pct)) - Inches(0.2)
    # Dot on timeline
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 mx + Inches(0.1), bar_y - Inches(0.12),
                                 Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GOLD
    dot.line.fill.background()
    dot.shadow.inherit = False
    # Date above
    add_textbox(slide, mx - Inches(0.3), bar_y - Inches(1.1),
                Inches(1.1), Inches(0.8),
                date, font_size=12, bold=True, color=NAVY,
                alignment=PP_ALIGN.CENTER)
    # Description below
    add_textbox(slide, mx - Inches(0.4), bar_y + Inches(0.35),
                Inches(1.4), Inches(1.0),
                desc, font_size=11, color=DGRAY,
                alignment=PP_ALIGN.CENTER)

# Phase boxes at bottom
phases = [
    ("Phase 1: Design & Prototype", "Mar \u2013 May 2026", NAVY),
    ("Phase 2: Regulatory & Security", "Jun \u2013 Dec 2026", LTNAVY),
    ("Phase 3: Scale Production", "Jan \u2013 May 2027", RGBColor(0x3D, 0x55, 0x8C)),
]

for i, (phase, period, color) in enumerate(phases):
    px = Inches(0.6) + i * Inches(4.2)
    py = Inches(5.0)
    add_rounded_rect(slide, px, py, Inches(3.9), Inches(1.2),
                     fill_color=color, line_color=GOLD, line_width=Pt(1))
    add_textbox(slide, px + Inches(0.1), py + Inches(0.15),
                Inches(3.7), Inches(0.5),
                phase, font_size=13, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, px + Inches(0.1), py + Inches(0.6),
                Inches(3.7), Inches(0.4),
                period, font_size=11, color=RGBColor(0xC0, 0xCC, 0xDD),
                alignment=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 14  --  Investment & Pricing
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, WHITE)
add_header_bar(slide, "Investment & Pricing")
add_footer(slide, 14)

# Key numbers
key_stats = [
    ("~$506K", "Phase 1-2\nDevelopment Investment"),
    ("$85", "Per-Unit Software\nCost at Scale"),
    ("30", "Units for Sword\n2026 Demo"),
    ("15,000/yr", "Full-Rate\nProduction Capacity"),
]

for i, (val, label) in enumerate(key_stats):
    sx = Inches(0.6) + i * Inches(3.15)
    add_stat_box(slide, sx, Inches(1.5), Inches(2.9), Inches(1.6),
                 val, label, num_color=GOLD, fill=NAVY,
                 label_color=RGBColor(0xC0, 0xCC, 0xDD))

# Value proposition bullets
value_bullets = [
    "Competitive Pricing: Per-unit cost at scale is competitive with existing "
    "hemodynamic monitoring systems \u2014 with significantly greater AI capability",
    "Prize Challenge Investment: Veteran Vectors' self-funded participation demonstrates "
    "commitment and de-risks the Government's investment",
    "Follow-On OTA Ready: Prototype contract-ready via DIU Other Transaction Authority "
    "pathway for rapid scaling",
    "Cost Avoidance: Early AI-assisted triage reduces downstream evacuation and treatment "
    "costs by improving initial casualty categorization accuracy",
    "Subscription Model Available: Optional SaaS model for cloud analytics tier reduces "
    "upfront acquisition cost",
]

bullet_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5),
               value_bullets, font_size=13, color=DGRAY,
               bold_first_word=True, spacing=Pt(8))


# ===========================================================================
# SLIDE 15  --  Call to Action / Contact
# ===========================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(slide, NAVY)

# Gold accent line
add_rect(slide, Inches(0.8), Inches(1.0), Inches(3), Inches(0.06), GOLD)

# Main headline
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.2),
            "VitalEdge AI", font_size=48, bold=True, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.7),
            "Saving Lives at the Point of Injury",
            font_size=26, bold=False, color=GOLD)

# Separator
add_rect(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.02),
         RGBColor(0x3D, 0x55, 0x8C))

# Contact details
contact_items = [
    "Veteran Vectors, LLC",
    "Service-Disabled Veteran-Owned Small Business (SDVOSB)",
    "",
    "Point of Contact:  [Name, Title]",
    "Email:  [email]@veteranvectors.com",
    "Phone:  [phone]",
    "Web:  www.veteranvectors.com",
]

ty = Inches(3.7)
for item in contact_items:
    if item == "":
        ty += Inches(0.15)
        continue
    add_textbox(slide, Inches(0.8), ty, Inches(6), Inches(0.4),
                item, font_size=14, bold=("Veteran Vectors, LLC" == item),
                color=WHITE if item == "Veteran Vectors, LLC" else RGBColor(0xC0, 0xCC, 0xDD))
    ty += Inches(0.35)

# Right side - key badges
badges = [
    "SDVOSB Certified",
    "PROJ00628",
    "Ready for Pitch Day",
    "Authentic Consortium",
]

for i, badge_text in enumerate(badges):
    bx = Inches(8.5)
    by = Inches(3.8) + i * Inches(0.85)
    add_rounded_rect(slide, bx, by, Inches(4.0), Inches(0.65),
                     fill_color=RGBColor(0x24, 0x3A, 0x5E),
                     line_color=GOLD, line_width=Pt(1))
    add_textbox(slide, bx + Inches(0.2), by + Inches(0.1), Inches(3.6), Inches(0.45),
                badge_text, font_size=14, bold=True, color=GOLD,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom classification
add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35),
         RGBColor(0x12, 0x1E, 0x36))
add_textbox(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35),
            "UNCLASSIFIED  |  Distribution Statement A: Approved for Public Release",
            font_size=9, color=RGBColor(0x70, 0x80, 0x99),
            alignment=PP_ALIGN.CENTER)


# ===========================================================================
# Save
# ===========================================================================
output_path = "/home/user/AC/DIU Proposal/03_Solution_Brief.pptx"
prs.save(output_path)
print(f"Solution brief saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
